# fitz_ai/ingestion/parser/plugins/lightweight.py
"""
Lightweight document parsers for PDF, DOCX, and PPTX.

Zero ML dependencies — uses pypdfium2, python-docx, and python-pptx
for basic text extraction. Included in base install.

For advanced layout analysis (tables, figures, reading order),
install docling: pip install fitz-ai[docs]
"""

from __future__ import annotations

import logging
import re
from dataclasses import dataclass, field
from typing import List, Set

from fitz_ai.core.document import DocumentElement, ElementType, ParsedDocument
from fitz_ai.ingestion.source.base import SourceFile

from .base_parser import BaseParser

logger = logging.getLogger(__name__)

# Extensions handled by lightweight parsers
LIGHTWEIGHT_EXTENSIONS: Set[str] = {
    ".pdf",
    ".docx",
    ".pptx",
}


@dataclass
class LightweightPDFParser(BaseParser):
    """
    Lightweight PDF parser using pypdfium2.

    Extracts text from digital PDFs (text-selectable). Does NOT handle
    scanned/image PDFs — use docling or a vision model for those.
    """

    plugin_name: str = field(default="lightweight_pdf")
    supported_extensions: Set[str] = field(default_factory=lambda: {".pdf"})

    def parse(self, file: SourceFile) -> ParsedDocument:
        try:
            import pypdfium2 as pdfium
        except ImportError:
            raise ImportError(
                "pypdfium2 required for PDF parsing. "
                "Install with: pip install pypdfium2"
            )

        file_bytes = self._read_file_bytes(file)
        pdf = pdfium.PdfDocument(file_bytes)

        elements: List[DocumentElement] = []
        for page_idx in range(len(pdf)):
            page = pdf[page_idx]
            text = page.get_textpage().get_text_range().strip()
            if text:
                # Try to extract headings from the text
                page_elements = self._extract_structure(text, page_idx + 1)
                elements.extend(page_elements)

        pdf.close()

        return ParsedDocument(
            source=file.uri,
            elements=elements,
            metadata=self._build_metadata(file, page_count=len(pdf) if pdf else 0),
        )

    def _extract_structure(self, text: str, page_num: int) -> List[DocumentElement]:
        """Extract basic structure from page text (headings vs paragraphs)."""
        elements = []
        paragraphs = re.split(r"\n{2,}", text)

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            # Heuristic: short lines in title case are likely headings
            lines = para.split("\n")
            if (
                len(lines) == 1
                and len(para) < 100
                and not para.endswith(".")
                and (para.isupper() or para.istitle())
            ):
                elements.append(
                    DocumentElement(
                        type=ElementType.HEADING,
                        content=para,
                        level=1,
                        metadata={"page": page_num},
                    )
                )
            else:
                elements.append(
                    DocumentElement(
                        type=ElementType.TEXT,
                        content=para,
                        metadata={"page": page_num},
                    )
                )

        return elements


@dataclass
class LightweightDOCXParser(BaseParser):
    """
    Lightweight DOCX parser using python-docx.

    Extracts text with heading structure from Word documents.
    """

    plugin_name: str = field(default="lightweight_docx")
    supported_extensions: Set[str] = field(default_factory=lambda: {".docx"})

    def parse(self, file: SourceFile) -> ParsedDocument:
        try:
            import docx
        except ImportError:
            raise ImportError(
                "python-docx required for DOCX parsing. "
                "Install with: pip install python-docx"
            )

        doc = docx.Document(str(file.local_path))
        elements: List[DocumentElement] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = (para.style.name or "").lower()

            if "heading" in style_name:
                # Extract heading level from style name (e.g., "Heading 2")
                level_match = re.search(r"\d+", style_name)
                level = int(level_match.group()) if level_match else 1
                elements.append(
                    DocumentElement(
                        type=ElementType.HEADING,
                        content=text,
                        level=level,
                    )
                )
            elif "list" in style_name:
                elements.append(
                    DocumentElement(
                        type=ElementType.LIST_ITEM,
                        content=text,
                    )
                )
            else:
                elements.append(
                    DocumentElement(
                        type=ElementType.TEXT,
                        content=text,
                    )
                )

        # Extract tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                elements.append(
                    DocumentElement(
                        type=ElementType.TABLE,
                        content="\n".join(rows),
                    )
                )

        return ParsedDocument(
            source=file.uri,
            elements=elements,
            metadata=self._build_metadata(file),
        )


@dataclass
class LightweightPPTXParser(BaseParser):
    """
    Lightweight PPTX parser using python-pptx.

    Extracts text from PowerPoint slides with slide-level structure.
    """

    plugin_name: str = field(default="lightweight_pptx")
    supported_extensions: Set[str] = field(default_factory=lambda: {".pptx"})

    def parse(self, file: SourceFile) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx required for PPTX parsing. "
                "Install with: pip install python-pptx"
            )

        prs = Presentation(str(file.local_path))
        elements: List[DocumentElement] = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts = []
            title = None

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if not text:
                        continue
                    if shape == slide.shapes.title:
                        title = text
                    else:
                        slide_texts.append(text)

            if title:
                elements.append(
                    DocumentElement(
                        type=ElementType.HEADING,
                        content=title,
                        level=1,
                        metadata={"slide": slide_idx},
                    )
                )

            for text in slide_texts:
                elements.append(
                    DocumentElement(
                        type=ElementType.TEXT,
                        content=text,
                        metadata={"slide": slide_idx},
                    )
                )

        return ParsedDocument(
            source=file.uri,
            elements=elements,
            metadata=self._build_metadata(file, slide_count=len(prs.slides)),
        )


__all__ = [
    "LightweightPDFParser",
    "LightweightDOCXParser",
    "LightweightPPTXParser",
    "LIGHTWEIGHT_EXTENSIONS",
]
